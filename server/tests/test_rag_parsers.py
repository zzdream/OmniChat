"""文档解析器单元测试。"""

import subprocess
from pathlib import Path

import pytest
import xlwt
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter

from app.services.rag.errors import IngestError
from app.services.rag.parsers import parse_document
from app.services.rag.parsers.excel import parse_xls_file, parse_xlsx_file
from app.services.rag.parsers.image import parse_image_file
from app.services.rag.parsers.pdf import parse_pdf_file
from app.services.rag.parsers.ppt import parse_pptx_file
from app.services.rag.parsers.text import parse_text_file
from app.services.rag.parsers.word import parse_doc_file, parse_docx_file


def _make_text_pdf(text: str, path: Path) -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.add_metadata({"/Title": text})
    with path.open("wb") as fp:
        writer.write(fp)


def test_parse_text_file_utf8(tmp_path: Path) -> None:
    file_path = tmp_path / "note.md"
    file_path.write_text("# Hello\n\n世界", encoding="utf-8")
    assert parse_text_file(file_path) == "# Hello\n\n世界"


def test_parse_text_file_unknown_encoding(tmp_path: Path) -> None:
    file_path = tmp_path / "bad.txt"
    file_path.write_bytes(b"\xff\xfe")
    with pytest.raises(IngestError, match="无法识别文件编码"):
        parse_text_file(file_path)


def test_parse_pdf_file_empty_text(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    file_path = tmp_path / "blank.pdf"
    _make_text_pdf("", file_path)
    monkeypatch.setattr(
        "app.services.rag.parsers.pdf.ocr_pdf_pages",
        lambda path, indices: {},
    )
    with pytest.raises(IngestError, match="未提取到文本"):
        parse_pdf_file(file_path)


def test_parse_pdf_file_uses_ocr_for_scan_pages(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    file_path = tmp_path / "scan.pdf"
    _make_text_pdf("", file_path)
    monkeypatch.setattr(
        "app.services.rag.parsers.pdf.ocr_pdf_pages",
        lambda path, indices: {1: "扫描页文字"},
    )
    assert parse_pdf_file(file_path) == "扫描页文字"


def test_parse_document_dispatches_by_extension(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    file_path = tmp_path / "demo.pdf"
    file_path.write_bytes(b"%PDF-1.4")

    monkeypatch.setattr(
        "app.services.rag.parsers.parse_pdf_file",
        lambda path: "PDF content",
    )
    assert parse_document(file_path) == "PDF content"


def test_parse_document_rejects_unsupported_extension(tmp_path: Path) -> None:
    file_path = tmp_path / "clip.gif"
    file_path.write_bytes(b"fake")
    with pytest.raises(IngestError, match="不支持的文件类型"):
        parse_document(file_path)


def _make_docx(path: Path, *paragraphs: str) -> None:
    document = Document()
    for paragraph in paragraphs:
        document.add_paragraph(paragraph)
    document.save(path)


def _make_doc_from_docx(docx_path: Path, doc_path: Path) -> None:
    subprocess.run(
        ["textutil", "-convert", "doc", "-output", str(doc_path), str(docx_path)],
        check=True,
        capture_output=True,
    )


def test_parse_docx_file_extracts_paragraphs_and_tables(tmp_path: Path) -> None:
    file_path = tmp_path / "report.docx"
    document = Document()
    document.add_paragraph("标题段落")
    table = document.add_table(rows=1, cols=2)
    table.rows[0].cells[0].text = "A1"
    table.rows[0].cells[1].text = "B1"
    document.save(file_path)

    text = parse_docx_file(file_path)
    assert "标题段落" in text
    assert "A1" in text and "B1" in text


def test_parse_docx_file_empty(tmp_path: Path) -> None:
    file_path = tmp_path / "empty.docx"
    _make_docx(file_path)
    with pytest.raises(IngestError, match="DOCX 未提取到文本"):
        parse_docx_file(file_path)


def test_parse_doc_file_extracts_text(tmp_path: Path) -> None:
    docx_path = tmp_path / "source.docx"
    doc_path = tmp_path / "legacy.doc"
    _make_docx(docx_path, "旧版 Word", "第二段")
    _make_doc_from_docx(docx_path, doc_path)

    text = parse_doc_file(doc_path)
    assert "旧版 Word" in text
    assert "第二段" in text


def test_parse_document_dispatches_word_extensions(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    docx_path = tmp_path / "memo.docx"
    doc_path = tmp_path / "memo.doc"
    docx_path.write_bytes(b"docx")
    doc_path.write_bytes(b"doc")

    monkeypatch.setattr(
        "app.services.rag.parsers.parse_docx_file",
        lambda path: "DOCX content",
    )
    monkeypatch.setattr(
        "app.services.rag.parsers.parse_doc_file",
        lambda path: "DOC content",
    )

    assert parse_document(docx_path) == "DOCX content"
    assert parse_document(doc_path) == "DOC content"


def _make_xlsx(path: Path) -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet["A1"] = "姓名"
    worksheet["B1"] = "分数"
    worksheet["A2"] = "张三"
    worksheet["B2"] = 95
    workbook.save(path)


def _make_xls(path: Path) -> None:
    workbook = xlwt.Workbook()
    sheet = workbook.add_sheet("Sheet1")
    sheet.write(0, 0, "产品")
    sheet.write(0, 1, "销量")
    sheet.write(1, 0, "A")
    sheet.write(1, 1, 100)
    workbook.save(path)


def _make_pptx(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "季度汇报"
    slide.placeholders[1].text = "收入增长 20%"
    presentation.save(path)


def test_parse_xlsx_file_extracts_rows(tmp_path: Path) -> None:
    file_path = tmp_path / "scores.xlsx"
    _make_xlsx(file_path)

    text = parse_xlsx_file(file_path)
    assert "姓名" in text and "张三" in text and "95" in text


def test_parse_xls_file_extracts_rows(tmp_path: Path) -> None:
    file_path = tmp_path / "sales.xls"
    _make_xls(file_path)

    text = parse_xls_file(file_path)
    assert "产品" in text and "销量" in text and "100" in text


def test_parse_pptx_file_extracts_slide_text(tmp_path: Path) -> None:
    file_path = tmp_path / "deck.pptx"
    _make_pptx(file_path)

    text = parse_pptx_file(file_path)
    assert "季度汇报" in text
    assert "收入增长 20%" in text


def test_parse_document_dispatches_excel_and_ppt_extensions(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    paths = {
        "book.xlsx": "XLSX content",
        "book.xls": "XLS content",
        "deck.pptx": "PPTX content",
        "deck.ppt": "PPT content",
    }
    for filename in paths:
        (tmp_path / filename).write_bytes(b"office")

    monkeypatch.setattr(
        "app.services.rag.parsers.parse_xlsx_file",
        lambda path: "XLSX content",
    )
    monkeypatch.setattr(
        "app.services.rag.parsers.parse_xls_file",
        lambda path: "XLS content",
    )
    monkeypatch.setattr(
        "app.services.rag.parsers.parse_pptx_file",
        lambda path: "PPTX content",
    )
    monkeypatch.setattr(
        "app.services.rag.parsers.parse_ppt_file",
        lambda path: "PPT content",
    )

    for filename, expected in paths.items():
        assert parse_document(tmp_path / filename) == expected


def test_parse_image_file_uses_ocr(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    file_path = tmp_path / "scan.png"
    file_path.write_bytes(b"\x89PNG\r\n")

    monkeypatch.setattr(
        "app.services.rag.parsers.image.ocr_image_path",
        lambda path: "图片中的文字",
    )
    assert parse_image_file(file_path) == "图片中的文字"


def test_parse_image_file_empty_ocr(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    file_path = tmp_path / "blank.png"
    file_path.write_bytes(b"\x89PNG\r\n")

    monkeypatch.setattr(
        "app.services.rag.parsers.image.ocr_image_path",
        lambda path: "",
    )
    with pytest.raises(IngestError, match="未识别到文字"):
        parse_image_file(file_path)
