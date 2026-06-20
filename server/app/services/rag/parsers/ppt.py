"""PowerPoint 文本提取：.pptx（python-pptx）、.ppt（系统转换）。"""

from pathlib import Path

from app.services.rag.errors import IngestError
from app.services.rag.parsers._utils import require_non_empty
from app.services.rag.parsers.legacy_office import extract_text_with_system_converter


def _extract_shape_text(shape: object) -> list[str]:
    if not getattr(shape, "has_text_frame", False):
        return []

    texts: list[str] = []
    text = getattr(shape, "text", "").strip()
    if text:
        texts.append(text)
        return texts

    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        line = paragraph.text.strip()
        if line:
            texts.append(line)
    return texts


def parse_pptx_file(file_path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise IngestError("PPTX 解析依赖未安装，请执行 pip install python-pptx") from exc

    try:
        presentation = Presentation(str(file_path))
    except Exception as exc:
        raise IngestError(f"无法读取 PPTX: {exc}") from exc

    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_texts: list[str] = []
        for shape in slide.shapes:
            slide_texts.extend(_extract_shape_text(shape))

        if slide_texts:
            parts.append(f"【第{index}页】\n" + "\n".join(slide_texts))

    return require_non_empty("\n\n".join(parts), empty_message="PPTX 未提取到文本")


def parse_ppt_file(file_path: Path) -> str:
    try:
        text = extract_text_with_system_converter(file_path)
    except IngestError:
        raise
    except Exception as exc:
        raise IngestError(f"无法读取 PPT: {exc}") from exc

    return require_non_empty(text, empty_message="PPT 未提取到文本")
